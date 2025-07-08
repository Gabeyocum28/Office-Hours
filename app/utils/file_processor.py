import os
import mimetypes
from pathlib import Path

def extract_text_from_file(file_path, file_type):
    """
    Extract text from various file types
    Returns extracted text or None if extraction fails
    """
    try:
        if file_type == 'text':
            return extract_from_txt(file_path)
        elif file_type == 'pdf':
            return extract_from_pdf(file_path)
        elif file_type == 'document':
            return extract_from_doc(file_path)
        elif file_type == 'presentation':
            return extract_from_ppt(file_path)
        elif file_type == 'image':
            return extract_from_image(file_path)
        elif file_type == 'video':
            return extract_from_video(file_path)
        else:
            return None
    except Exception as e:
        print(f"Error extracting text from {file_path}: {e}")
        return None

def extract_from_txt(file_path):
    """Extract text from plain text files"""
    try:
        with open(file_path, 'r', encoding='utf-8') as file:
            return file.read()
    except UnicodeDecodeError:
        # Try with different encoding
        with open(file_path, 'r', encoding='latin-1') as file:
            return file.read()

def extract_from_pdf(file_path):
    """Extract text from PDF files using PyPDF2"""
    try:
        import PyPDF2
        
        text = ""
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
        return text.strip()
    except ImportError:
        print("PyPDF2 not installed. Install with: pip install PyPDF2")
        return None
    except Exception as e:
        print(f"PDF extraction error: {e}")
        return None

def extract_from_doc(file_path):
    """Extract text from Word documents using python-docx"""
    try:
        from docx import Document
        
        doc = Document(file_path)
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        return text.strip()
    except ImportError:
        print("python-docx not installed. Install with: pip install python-docx")
        return None
    except Exception as e:
        print(f"Word document extraction error: {e}")
        return None

def extract_from_ppt(file_path):
    """Extract text from PowerPoint presentations using python-pptx"""
    try:
        from pptx import Presentation
        
        prs = Presentation(file_path)
        text = ""
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text.strip()
    except ImportError:
        print("python-pptx not installed. Install with: pip install python-pptx")
        return None
    except Exception as e:
        print(f"PowerPoint extraction error: {e}")
        return None

def extract_from_image(file_path):
    """Extract text from images using OCR (Tesseract)"""
    try:
        from PIL import Image
        import pytesseract
        
        # Open image
        image = Image.open(file_path)
        
        # Extract text using OCR
        text = pytesseract.image_to_string(image)
        return text.strip()
    except ImportError:
        print("OCR dependencies not installed. Install with: pip install pillow pytesseract")
        print("Also install Tesseract OCR: https://github.com/tesseract-ocr/tesseract")
        return f"[IMAGE: {os.path.basename(file_path)} - OCR not available]"
    except Exception as e:
        print(f"Image OCR extraction error: {e}")
        return f"[IMAGE: {os.path.basename(file_path)} - text extraction failed]"

def extract_from_video(file_path):
    """Extract metadata/transcript from video files"""
    # For now, just return a placeholder
    # In the future, you could integrate with:
    # - Speech-to-text APIs (OpenAI Whisper, Google Speech-to-Text)
    # - Video frame analysis for slides/text
    filename = os.path.basename(file_path)
    return f"[VIDEO: {filename} - transcript extraction not yet implemented]"

def get_file_summary(file_path, extracted_text):
    """
    Generate a brief summary of the file content
    This could be used for quick file identification
    """
    if not extracted_text:
        return f"File: {os.path.basename(file_path)} (no text extracted)"
    
    # Simple summary - first 200 characters
    summary = extracted_text[:200]
    if len(extracted_text) > 200:
        summary += "..."
    
    return f"File: {os.path.basename(file_path)}\nContent preview: {summary}"